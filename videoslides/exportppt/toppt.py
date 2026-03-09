import os, time
import cv2
import numpy as np
from scipy import signal
from pptx import Presentation
from pptx.util import Inches
from io import BytesIO  # or from io import StringIO for text-based images
from PIL import Image
import argparse # Added for command-line arguments
import concurrent
from concurrent.futures import ThreadPoolExecutor # Added for multi-threading


def recursive_list_file_paths(PATH, extension=[]):
    '''
    PATH : string , path to search
    extension : list of strings, file extension(s) to included
    '''
    result = [os.path.join(dp, f) for dp, dn, filenames in os.walk(PATH) for f in filenames if os.path.splitext(f)[1] in extension] if len(extension)>0 else \
    [os.path.join(dp, f) for dp, dn, filenames in os.walk(PATH) for f in filenames]
    return result

def moving_mean_std(array, d=10):
    length=array.shape[0]
    if length<d:
        print("array too short")
        return None
    idx = np.arange(length-d+1).reshape(-1,1) + np.arange(d).reshape(1,-1)
    mean_array=array[idx].mean(axis=1)
    std_array=array[idx].std(axis=1)
    return mean_array, std_array

def select_frames(all_frame_changes, height=None, distance=6):

    length = all_frame_changes.shape[0]
    if height is None:
        mean, std = moving_mean_std(all_frame_changes)
        pad_size= length-mean.shape[0]
        padding=(pad_size-(pad_size//2), pad_size//2)
        height =np.pad(mean, padding, 'constant', constant_values=0) + np.pad(std, padding, 'constant', constant_values=0)
    #height= np.percentile(all_frame_changes[:,1], pc)
    frame_peaks, _ = signal.find_peaks(all_frame_changes,height=height, distance=distance)
    # pixel_peaks, _ = signal.find_peaks(all_frame_changes, distance=2)
    pixel_valley, _ = signal.find_peaks(-all_frame_changes, distance=2)
    sections = frame_peaks.shape[0] # not include last section
    selection_frames = []
    
    if sections < 1 : # when no sharp change
        selection_frames.append(all_frame_changes.shape[0]//2)
        return selection_frames

    # from beginning to 1st sharp change
    # pixel_peaks_in_section = pixel_peaks[ (pixel_peaks < frame_peaks[0]) ]
    pixel_valley_insection = pixel_valley[ (pixel_valley < frame_peaks[0]) ]
    
    f = pixel_valley_insection[-1] if pixel_valley_insection.size != 0 else (frame_peaks[0])//2
    selection_frames.append(f) 
    
    for i in range(1,sections):
        pixel_valley_in_section = pixel_valley[ (pixel_valley > frame_peaks[i-1]) & (pixel_valley < frame_peaks[i]) ]
        f = (frame_peaks[i]-3) if pixel_valley_in_section.size == 0 else pixel_valley_in_section[-1]
        selection_frames.append(f)
    # after last sharp change
    pixel_valley_insection = pixel_valley[ (pixel_valley > frame_peaks[-1]) ]
    f = pixel_valley_insection[-1] if pixel_valley_insection.size != 0 else (frame_peaks[-1]+all_frame_changes.shape[0]-1)//2
    selection_frames.append(f)
    return selection_frames

def create_powerpoint_with_select_frames(video_path, selection_frames, output_pptx_name="presentation.pptx"):
    """
    Creates a PowerPoint presentation and inserts images into slides.
    Args:
        selection_frames: A list of frame numbers you want to insert.
        output_filename: The name of the PowerPoint file to create.
    """
    # Create a PowerPoint presentation object
    prs = Presentation()
    prs.slide_width = Inches(7.5*16/9)
    prs.slide_height = Inches(7.5)
    
    blank_slide_layout = prs.slide_layouts[6]  # Index 6 usually corresponds to a blank slide

    # Open the video file
    video = cv2.VideoCapture(video_path)

    # Iterate over selected frame in the video
    for f in selection_frames:
        video.set(cv2.CAP_PROP_POS_FRAMES, f)
        success, frame = video.read()
        if success:
            frame_rgb = cv2.cvtColor(frame, cv2.COLOR_BGR2RGB)
            image = Image.fromarray(frame_rgb)
            # Create an in-memory byte stream to save to memory
            image_stream = BytesIO()
            image.save(image_stream, format='PNG')
            image_stream.seek(0)  # Reset the stream position to the beginning

            slide = prs.slides.add_slide(blank_slide_layout)

            # Set image position and size
            left, top = Inches(0), Inches(0)
            width, height = Inches(7.5*16/9), Inches(7.5)

            slide.shapes.add_picture(image_stream, left, top, width, height)
            
    prs.save(output_pptx_name)
    print(f"PowerPoint presentation saved to {output_pptx_name}")

    # Release the video capture object
    video.release()

# def video_to_pptx_pipeline(video_path):
#     print(video_path)
#     export_dir = os.path.splitext(video_path)[0]
#     output_pptx_name = export_dir + '.pptx'

#     # analyze video
#     t0=time.time()
#     all_frame_changes = analyze_frames(video_path)
#     all_frame_changes = np.array(all_frame_changes, dtype=np.float32)
#     t1=time.time()
#     print('analyze video: ', t1-t0)
    
#     # select frame
#     t0=time.time()
#     FRAME_RATE = 30
#     distance = FRAME_RATE//5
#     height = all_frame_changes[:,1].mean()+all_frame_changes[:,1].std()
#     selection_frames = select_frames(all_frame_changes, height=height, distance=distance)
#     t1=time.time()
#     print('select frame: ', t1-t0)
#     # v2.0
#     # frame to pptx
#     t0=time.time()
#     create_powerpoint_with_select_frames(video_path, selection_frames, output_pptx_name=output_pptx_name)
    
#     t1=time.time()
#     print('frames to pptx: ',t1-t0)

def main():
    parser = argparse.ArgumentParser(description="Analyze videos and generate PowerPoint presentations from selected frames.")
    parser.add_argument("--video_folder", type=str, required=True,
                        help="Path to the folder containing video files.")
    parser.add_argument("--video_extension", type=str, nargs='+', default=['.mp4', '.MP4'],
                        help="List of video file extensions to include (e.g., .mp4 .avi).")
    parser.add_argument("--num_of_thread", type=int, default=2,
                        help="Number of threads to use for parallel video processing.")
    
    args = parser.parse_args()

    video_folder = args.video_folder
    # Normalize extensions to lowercase for consistent matching
    video_extension = [ext.lower() for ext in args.video_extension] 
    num_of_thread = args.num_of_thread

    print(f"Searching for videos in: {video_folder}")
    print(f"Including extensions: {video_extension}")
    print(f"Using {num_of_thread} threads.")

    video_list = recursive_list_file_paths(
        video_folder,
        extension=video_extension)

    if not video_list:
        print(f"No video files found with extensions {video_extension} in {video_folder}")
        return

    print(f"Found {len(video_list)} video files.")

    # Multi-threading the video_to_pptx_pipeline
    with ThreadPoolExecutor(max_workers=num_of_thread) as executor:
        # Submit each video processing task to the executor
        futures = {executor.submit(video_to_pptx_pipeline, v): v for v in video_list}
        
        # Wait for tasks to complete and handle results/exceptions
        for future in concurrent.futures.as_completed(futures):
            video_path = futures[future]
            try:
                future.result() # This will re-raise any exception that occurred in the thread
            except Exception as exc:
                print(f'Error processing {video_path}: {exc}')
            else:
                print(f'Finished processing {video_path}.')

if __name__ == '__main__':
    main()
